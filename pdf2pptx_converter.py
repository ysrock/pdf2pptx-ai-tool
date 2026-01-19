import sys
import os
import pdfplumber
import numpy as np
import cv2
from rapidocr_onnxruntime import RapidOCR
from pptx import Presentation
from pptx.util import Pt
from simple_lama_inpainting import SimpleLama
from PIL import Image
import torch
import gc
import tempfile
import shutil
import threading
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk

# --- Global Config ---
DEFAULT_DPI = 250

# --- Conversion Logic (Refactored for GUI) ---
class ConverterLogic:
    def __init__(self, log_callback=print):
        self.log = log_callback
        self.stop_flag = False
        self.ocr_engine = None
        self.lama = None
        self.is_gpu = False
        self._initialized = False
    
    def initialize_models(self):
        """Initialize OCR and LaMa models once for the entire batch"""
        if self._initialized:
            return True
            
        # Initialize OCR
        try:
            self.log("Initializing OCR engine...")
            self.ocr_engine = RapidOCR()
        except Exception as e:
            self.log(f"Error initializing OCR: {e}")
            return False

        # Initialize LaMa
        try:
            self.log("Initializing LaMa AI model...")
            self.is_gpu = False
            if torch.cuda.is_available():
                self.log(f"✅ GPU Detected: {torch.cuda.get_device_name(0)}")
                self.is_gpu = True
            else:
                self.log("⚠️ GPU NOT Detected. Running on CPU (Slow Mode)")

            # Bundle logic
            if getattr(sys, 'frozen', False):
                 model_path = os.path.join(sys._MEIPASS, 'big-lama.pt')
                 self.log(f"Loading bundled model...")
                 fake_home = os.path.join(sys._MEIPASS, 'torch_cache')
                 if not os.path.exists(fake_home): os.makedirs(fake_home)
                 checkpoint_dir = os.path.join(fake_home, 'hub', 'checkpoints')
                 if not os.path.exists(checkpoint_dir): os.makedirs(checkpoint_dir)
                 target_model = os.path.join(checkpoint_dir, 'big-lama.pt')
                 if not os.path.exists(target_model):
                     shutil.copy2(model_path, target_model)
                 os.environ['TORCH_HOME'] = fake_home
            
            # Force CPU for SimpleLama load if no GPU is available or if we are in CPU mode
            if not self.is_gpu:
                 # Monkeypatch torch.jit.load (NOT torch.load) because SimpleLama uses JIT
                 original_jit_load = torch.jit.load
                 def safe_jit_load(*args, **kwargs):
                     if 'map_location' not in kwargs:
                         kwargs['map_location'] = torch.device('cpu')
                     return original_jit_load(*args, **kwargs)
                 torch.jit.load = safe_jit_load
                 
                 self.lama = SimpleLama()
                 
                 # Restore original load
                 torch.jit.load = original_jit_load
            else:
                 self.lama = SimpleLama()
                 
            self.log("LaMa model loaded.")
            self._initialized = True
            return True
        except Exception as e:
            self.log(f"Error initializing LaMa: {e}")
            # If LaMa fails, we can still try to proceed without it, or return False
            # For now, let's return False to indicate a critical failure.
            return False

    def cleanup_file(self):
        """Aggressive memory cleanup after processing each file"""
        # Force Python garbage collection multiple times
        for _ in range(3):
            gc.collect()
        
        # Clear CUDA cache if using GPU
        if self.is_gpu and torch.cuda.is_available():
            torch.cuda.empty_cache()
            torch.cuda.synchronize()
        
        # Additional cleanup for CPU memory
        if hasattr(torch, '_C') and hasattr(torch._C, '_cuda_clearCublasWorkspaces'):
            try:
                torch._C._cuda_clearCublasWorkspaces()
            except:
                pass

    def convert(self, pdf_path, pptx_path, dpi, dilation_size, debug_mode=False):
        self.log(f"Starting Conversion...")
        self.log(f"Input: {pdf_path}")
        self.log(f"Output: {pptx_path}")
        self.log(f"Quality: {dpi} DPI")
        self.log(f"Cleaner Strength: {dilation_size}")
        self.log(f"Debug Mode: {'ON' if debug_mode else 'OFF'}")
        
        # Ensure models are initialized
        if not self._initialized:
            if not self.initialize_models():
                return False

        try:
            pdf_file = pdfplumber.open(pdf_path)
            total_pages = len(pdf_file.pages)
            self.log(f"Total pages: {total_pages}")
        except Exception as e:
            self.log(f"Error opening PDF: {e}")
            return False

        prs = Presentation()
        
        with tempfile.TemporaryDirectory() as temp_dir:
            for i, page in enumerate(pdf_file.pages):
                if self.stop_flag:
                    self.log("Conversion Stopped by User.")
                    return False
                    
                self.log(f"Processing page {i + 1}/{total_pages}...")
                
                try:
                    pdf_width_pt = page.width
                    pdf_height_pt = page.height
                    
                    if i == 0:
                        prs.slide_width = Pt(pdf_width_pt)
                        prs.slide_height = Pt(pdf_height_pt)
                    
                    blank_slide_layout = prs.slide_layouts[6] 
                    slide = prs.slides.add_slide(blank_slide_layout)

                    page_image_obj = page.to_image(resolution=dpi)
                    pil_image = page_image_obj.original
                    img_np = np.array(pil_image)
                    
                    mask = np.zeros(img_np.shape[:2], dtype=np.uint8)
                    ocr_result, _ = self.ocr_engine(img_np)
                    
                    text_blocks = []
                    box_heights = []
                    scale_factor = 72 / dpi
                    
                    if ocr_result:
                        for item in ocr_result:
                            box, text, score = item
                            
                            # Expand the polygon slightly (5-10%) to cover edges/punctuation
                            poly = np.array(box, dtype=np.float32)
                            center = np.mean(poly, axis=0)
                            # Expansion vector: (Point - Center) * Scale + Center
                            expanded_poly = center + (poly - center) * 1.10
                            
                            pts = expanded_poly.astype(np.int32).reshape((-1, 1, 2))
                            cv2.fillPoly(mask, [pts], 255)
                            
                            xs = [p[0] for p in box]
                            ys = [p[1] for p in box]
                            text_blocks.append({
                                'text': text,
                                'x': min(xs) * scale_factor, 
                                'y': min(ys) * scale_factor, 
                                'w': (max(xs) - min(xs)) * scale_factor, 
                                'h': (max(ys) - min(ys)) * scale_factor
                            })
                            
                            # Collect height for adaptive dilation
                            h_px = max(ys) - min(ys)
                            box_heights.append(h_px)
                        
                        # Adaptive Dilation Calculation
                        # Instead of a fixed global kernel, we compute it based on the text size.
                        # Logic: Larger text needs more "bleed" to erase effectively.
                        if box_heights:
                            avg_height = sum(box_heights) / len(box_heights)
                            # Heuristic: 20% of text height, clamped between 15 and 40
                            adaptive_dil = int(avg_height * 0.20)
                            adaptive_dil = max(15, min(adaptive_dil, 40))
                            
                            # If manual override is significantly different, maybe use that? 
                            # But user said "I don't want to tune". So let's use the Dynamic value 
                            # unless the user explicitly changed the slider from default in a future version.
                            # For now, we mix them: Max(Manual, Adaptive) to be safe.
                            final_dil = max(dilation_size, adaptive_dil)
                        else:
                            final_dil = dilation_size

                        # Optimization: Morphological CLOSE to connect nearby characters
                        kernel_close = np.ones((5, 5), np.uint8)
                        mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, kernel_close)

                        # Apply Final Dilation
                        kernel = np.ones((final_dil, final_dil), np.uint8)
                        mask = cv2.dilate(mask, kernel, iterations=1)
                    
                    mask_pil = Image.fromarray(mask)
                    with torch.no_grad():
                        cleaned_image_pil = self.lama(pil_image, mask_pil)
                    
                    bg_image_path = os.path.join(temp_dir, f"bg_page_{i}.png")
                    cleaned_image_pil.save(bg_image_path)
                    
                    # Debug Mode: Export Clean Background
                    if debug_mode:
                        debug_dir = pptx_path.replace(".pptx", "_debug_images")
                        if not os.path.exists(debug_dir):
                            os.makedirs(debug_dir, exist_ok=True)
                        debug_save_path = os.path.join(debug_dir, f"page_{i+1}_clean_bg.png")
                        cleaned_image_pil.save(debug_save_path)
                        
                    slide.shapes.add_picture(bg_image_path, Pt(0), Pt(0), width=Pt(pdf_width_pt), height=Pt(pdf_height_pt))
                    
                    for block in text_blocks:
                        try:
                            txBox = slide.shapes.add_textbox(Pt(block['x']), Pt(block['y']), Pt(block['w']), Pt(block['h']))
                            tf = txBox.text_frame
                            tf.word_wrap = True
                            tf.text = block['text']
                            if block['h'] > 0:
                                tf.paragraphs[0].font.size = Pt(block['h'] * 0.8)
                        except:
                            pass
                    
                    del page_image_obj, pil_image, img_np, mask, mask_pil, cleaned_image_pil
                    if self.is_gpu: torch.cuda.empty_cache()
                    gc.collect()
                    
                except Exception as e:
                    self.log(f"Error processing page {i}: {e}")

            pdf_file.close()
            
            self.log(f"Saving to: {pptx_path}")
            try:
                prs.save(pptx_path)
                self.log("Conversion Success!")
                # Explicitly delete the Presentation object to free memory
                del prs
                gc.collect()
                return True
            except Exception as e:
                self.log(f"Error saving PPTX: {e}")
                del prs
                gc.collect()
                return False

# --- GUI Application (Modern Content) ---
# Try to import tkinterdnd2 for native drag-drop support
try:
    from tkinterdnd2 import DND_FILES, TkinterDnD
    HAS_DND = True
except ImportError:
    HAS_DND = False

class App:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF to PPTX Converter AI")
        self.root.geometry("900x650")
        
        # --- Theme Config ---
        self.colors = {
            "bg_main": "#2C2C2C",
            "bg_sec":  "#383838",
            "fg_text": "#FFFFFF",
            "accent":  "#3B8ED0",
            "accent_hover": "#3071A9",
            "success": "#28A745",
            "error":   "#DC3545",
            "entry_bg": "#454545"
        }
        
        self.setup_style()
        self.root.configure(bg=self.colors["bg_main"])
        
        # --- File Queue ---
        self.file_queue = []
        
        # --- Main Layout ---
        # Header
        self.create_header()
        
        # Content Area (Split Left/Right)
        content_frame = ttk.Frame(root, style="Main.TFrame")
        content_frame.pack(fill="both", expand=True, padx=20, pady=20)
        
        # Left Panel: File Queue
        self.create_left_panel(content_frame)
        
        # Right Panel: Settings & Actions
        self.create_right_panel(content_frame)
        
        # Footer: Logs
        self.create_footer()
        
        if HAS_DND:
            self.log("Ready. Drag PDF files to the list on the left.")
        else:
            self.log("Ready. Click '+ Add PDFs' to start.")

    def setup_style(self):
        style = ttk.Style()
        style.theme_use('clam')  # 'clam' allows easier color customization than 'vista'
        
        # Common formatting
        style.configure(".", background=self.colors["bg_main"], foreground=self.colors["fg_text"], font=("Segoe UI", 10))
        style.configure("TFrame", background=self.colors["bg_main"])
        style.configure("Main.TFrame", background=self.colors["bg_main"])
        
        # Label
        style.configure("TLabel", background=self.colors["bg_main"], foreground=self.colors["fg_text"], font=("Segoe UI", 10))
        style.configure("Header.TLabel", font=("Segoe UI", 18, "bold"), foreground=self.colors["fg_text"])
        style.configure("SubHeader.TLabel", font=("Segoe UI", 12, "bold"), foreground="#AAAAAA")
        style.configure("Card.TLabel", background=self.colors["bg_sec"], font=("Segoe UI", 10, "bold"))
        
        # Button ("Accent" style)
        style.configure("Accent.TButton", 
                        background=self.colors["accent"], 
                        foreground="white", 
                        font=("Segoe UI", 10, "bold"),
                        borderwidth=0, focuscolor=self.colors["bg_main"])
        style.map("Accent.TButton", background=[("active", self.colors["accent_hover"])])
        
        # Button ("Secondary" style)
        style.configure("Sec.TButton", 
                        background=self.colors["bg_sec"], 
                        foreground="white", 
                        borderwidth=0, focuscolor=self.colors["bg_main"])
        style.map("Sec.TButton", background=[("active", "#505050")])
        
        # Frame "Card" background
        style.configure("Card.TFrame", background=self.colors["bg_sec"], relief="flat")
        
        # LabelFrame
        style.configure("TLabelframe", background=self.colors["bg_sec"], foreground=self.colors["fg_text"], bordercolor="#555555")
        style.configure("TLabelframe.Label", background=self.colors["bg_sec"], foreground=self.colors["fg_text"])
        
        # Checkbutton
        style.configure("TCheckbutton", background=self.colors["bg_sec"], foreground=self.colors["fg_text"], indicatorbackground=self.colors["bg_sec"], selectcolor=self.colors["accent"])
        
        # Spinbox - Fix text visibility
        style.configure("TSpinbox", 
                        fieldbackground=self.colors["entry_bg"], 
                        foreground=self.colors["fg_text"],
                        arrowcolor=self.colors["fg_text"])

    def create_header(self):
        header_frame = ttk.Frame(self.root, padding="20 15 20 0")
        header_frame.pack(fill="x")
        
        lbl_title = ttk.Label(header_frame, text="PDF to PPTX Converter", style="Header.TLabel")
        lbl_title.pack(anchor="w")
        
        lbl_sub = ttk.Label(header_frame, text="AI-Powered Layout Reconstruction", style="SubHeader.TLabel")
        lbl_sub.pack(anchor="w")

    def create_left_panel(self, parent):
        frame = ttk.Frame(parent, padding="0 0 10 0")
        frame.pack(side="left", fill="both", expand=True)
        
        # Card container
        card = ttk.Frame(frame, style="Card.TFrame", padding=15)
        card.pack(fill="both", expand=True)
        
        ttk.Label(card, text="File Queue", style="Card.TLabel", background=self.colors["bg_sec"]).pack(anchor="w", pady=(0, 10))
        
        # Listbox with custom colors
        self.listbox = tk.Listbox(card, 
                                  bg=self.colors["entry_bg"], 
                                  fg="white", 
                                  selectbackground=self.colors["accent"],
                                  selectforeground="white",
                                  borderwidth=0, 
                                  highlightthickness=1,
                                  highlightbackground="#555555",
                                  font=("Segoe UI", 10),
                                  height=10)
        self.listbox.pack(side="left", fill="both", expand=True)
        
        scrollbar = ttk.Scrollbar(card, orient="vertical", command=self.listbox.yview)
        scrollbar.pack(side="right", fill="y")
        self.listbox.config(yscrollcommand=scrollbar.set)
        
        # Drag Drop
        if HAS_DND:
            self.listbox.drop_target_register(DND_FILES)
            self.listbox.dnd_bind('<<Drop>>', self.on_drop)

    def create_right_panel(self, parent):
        frame = ttk.Frame(parent, padding="10 0 0 0")
        frame.pack(side="right", fill="y")
        
        # --- Actions Card ---
        card_actions = ttk.Frame(frame, style="Card.TFrame", padding=15)
        card_actions.pack(fill="x", pady=(0, 15))
        
        ttk.Label(card_actions, text="Manage Queue", style="Card.TLabel", background=self.colors["bg_sec"]).pack(anchor="w", pady=(0,10))
        
        btn_add = ttk.Button(card_actions, text="+ Add PDFs", command=self.browse_files, style="Accent.TButton", cursor="hand2")
        btn_add.pack(fill="x", pady=2)
        
        btn_del = ttk.Button(card_actions, text="Remove Selected", command=self.remove_selected, style="Sec.TButton", cursor="hand2")
        btn_del.pack(fill="x", pady=2)
        
        btn_clr = ttk.Button(card_actions, text="Clear All", command=self.clear_queue, style="Sec.TButton", cursor="hand2")
        btn_clr.pack(fill="x", pady=2)
        
        # --- Settings Card ---
        card_settings = ttk.Frame(frame, style="Card.TFrame", padding=15)
        card_settings.pack(fill="x", pady=(0, 15))
        
        ttk.Label(card_settings, text="Conversion Settings", style="Card.TLabel", background=self.colors["bg_sec"]).pack(anchor="w", pady=(0,10))
        
        # DPI
        row1 = ttk.Frame(card_settings, style="Card.TFrame")
        row1.pack(fill="x", pady=2)
        ttk.Label(row1, text="Quality (DPI):", background=self.colors["bg_sec"]).pack(side="left")
        self.dpi_var = tk.IntVar(value=100)
        spin_dpi = ttk.Spinbox(row1, from_=72, to=300, textvariable=self.dpi_var, width=5)
        spin_dpi.pack(side="right")
        
        # Dilation
        row2 = ttk.Frame(card_settings, style="Card.TFrame")
        row2.pack(fill="x", pady=5)
        ttk.Label(row2, text="Artifact Cleanup:", background=self.colors["bg_sec"]).pack(side="left")
        self.dil_var = tk.IntVar(value=15)
        spin_dil = ttk.Spinbox(row2, from_=0, to=50, textvariable=self.dil_var, width=5)
        spin_dil.pack(side="right")
        
        # Debug
        self.debug_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(card_settings, text="Export Debug Images", variable=self.debug_var, style="TCheckbutton").pack(anchor="w", pady=5)
        
        # --- Start Button ---
        self.btn_convert = ttk.Button(frame, text="START BATCH CONVERSION", command=self.start_batch, style="Accent.TButton", cursor="hand2")
        self.btn_convert.pack(fill="x", pady=10, ipady=5)

    def create_footer(self):
        footer_frame = ttk.Frame(self.root, padding="20 0 20 20")
        footer_frame.pack(fill="both")
        
        lbl_log = ttk.Label(footer_frame, text="Process Log", style="SubHeader.TLabel", font=("Segoe UI", 10, "bold"))
        lbl_log.pack(anchor="w", pady=(0, 5))
        
        self.text_log = scrolledtext.ScrolledText(footer_frame, 
                                                  height=8, 
                                                  bg=self.colors["bg_sec"], 
                                                  fg="#CCCCCC",
                                                  font=("Consolas", 9),
                                                  borderwidth=0,
                                                  state="normal")
        self.text_log.pack(fill="both", expand=True)

    # --- Logic Methods (Same as before) ---
    def log(self, msg):
        self.text_log.insert(tk.END, f"> {msg}\n")
        self.text_log.see(tk.END)
        self.root.update_idletasks()

    def on_drop(self, event):
        data = event.data
        if '{' in data:
            import re
            files = re.findall(r'\{([^}]+)\}', data)
            data_clean = re.sub(r'\{[^}]+\}', '', data)
            files.extend(data_clean.split())
        else:
            files = data.split()
        
        count = 0
        for f in files:
            f = f.strip()
            if f.lower().endswith('.pdf') and f not in self.file_queue:
                self.file_queue.append(f)
                self.listbox.insert(tk.END, os.path.basename(f))
                count += 1
        
        if count > 0:
            self.log(f"Queue updated: +{count} files (Total: {len(self.file_queue)})")

    def browse_files(self):
        files = filedialog.askopenfilenames(filetypes=[("PDF Files", "*.pdf")])
        count = 0
        for f in files:
            if f not in self.file_queue:
                self.file_queue.append(f)
                self.listbox.insert(tk.END, os.path.basename(f))
                count += 1
        if count > 0:
            self.log(f"Queue updated: +{count} files (Total: {len(self.file_queue)})")

    def remove_selected(self):
        selected = list(self.listbox.curselection())
        selected.reverse()
        for idx in selected:
            del self.file_queue[idx]
            self.listbox.delete(idx)
        if selected:
            self.log(f"Removed {len(selected)} files.")

    def clear_queue(self):
        self.file_queue.clear()
        self.listbox.delete(0, tk.END)
        self.log("Queue cleared.")

    def start_batch(self):
        if not self.file_queue:
            messagebox.showerror("Error", "No PDF files in queue.")
            return
        
        self.btn_convert.config(state="disabled", text="PROCESSING...")
        dpi = self.dpi_var.get()
        dil = self.dil_var.get()
        debug = self.debug_var.get()
        
        files_to_process = list(self.file_queue)
        
        t = threading.Thread(target=self.run_batch, args=(files_to_process, dpi, dil, debug))
        t.daemon = True
        t.start()

    def run_batch(self, files, dpi, dil, debug):
        total = len(files)
        success_count = 0
        
        self.log_thread_safe(f"--- STARTING BATCH: {total} FILES ---")
        
        converter = ConverterLogic(log_callback=self.log_thread_safe)
        
        self.log_thread_safe("Initializing AI Models...")
        if not converter.initialize_models():
            self.log_thread_safe("CRITICAL: Failed to init models.")
            self.reset_ui()
            return
        
        for i, pdf_path in enumerate(files):
            self.log_thread_safe(f"[{i+1}/{total}] Processing: {os.path.basename(pdf_path)}")
            
            base_name = os.path.splitext(pdf_path)[0]
            pptx_path = f"{base_name}_Editable.pptx"
            
            try:
                result = converter.convert(pdf_path, pptx_path, dpi, dil, debug)
                if result:
                    success_count += 1
            except Exception as e:
                self.log_thread_safe(f"Error: {e}")
            
            converter.cleanup_file()
        
        self.log_thread_safe(f"--- COMPLETED: {success_count}/{total} SUCCEEDED ---")
        self.reset_ui()

    def reset_ui(self):
        self.root.after(0, lambda: self.btn_convert.config(state="normal", text="START BATCH CONVERSION"))

    def log_thread_safe(self, msg):
        self.root.after(0, lambda: self.log(msg))

if __name__ == "__main__":
    # If args provided, run headless
    if len(sys.argv) > 1:
        pdf_file = sys.argv[1]
        pptx_file = sys.argv[2] if len(sys.argv) > 2 else os.path.splitext(pdf_file)[0] + "_Editable.pptx"
        dpi = int(sys.argv[3]) if len(sys.argv) > 3 else 100
        dilation = int(sys.argv[4]) if len(sys.argv) > 4 else 15
        
        logic = ConverterLogic()
        if os.path.exists(pdf_file):
            print(f"CLI Mode: DPI={dpi}, Dilation={dilation}")
            logic.convert(pdf_file, pptx_file, dpi, dilation)
    else:
        # GUI Mode
        if HAS_DND:
            root = TkinterDnD.Tk()
        else:
            root = tk.Tk()
        app = App(root)
        root.mainloop()

